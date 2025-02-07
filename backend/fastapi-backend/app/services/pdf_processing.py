from typing import Tuple, Any
import pdfplumber
from PIL import Image
import matplotlib.pyplot as plt
from pydub import AudioSegment
from pptx import Presentation
import os

def extract_text(pdf_path: str) -> str:
    with pdfplumber.open(pdf_path) as pdf:
        text = ''
        for page in pdf.pages:
            text += page.extract_text() + '\n'
    return text.strip()

def generate_graphical_abstract(text: str, output_path: str) -> None:
    plt.figure(figsize=(10, 5))
    plt.text(0.5, 0.5, text, fontsize=12, ha='center', va='center')
    plt.axis('off')
    plt.savefig(output_path, bbox_inches='tight')
    plt.close()

def generate_short_content(text: str, output_dir: str) -> None:
    # Generate PPT
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    textbox = slide.shapes.add_textbox(left=0, top=0, width=ppt.slide_width, height=ppt.slide_height)
    textbox.text = text
    ppt.save(os.path.join(output_dir, 'presentation.pptx'))

    # Generate audio
    audio = AudioSegment.silent(duration=1000)  # Placeholder for audio generation
    audio.export(os.path.join(output_dir, 'audio.mp3'), format='mp3')